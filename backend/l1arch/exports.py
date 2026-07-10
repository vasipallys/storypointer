"""Export the L1 executive summary to Markdown, Word and PowerPoint.

Reuses the Markdown/Mermaid parsing + rendering helpers from the requirements
exporter so the "living document" can also leave the app as a deck or doc — the
Markdown+Mermaid view remains the source of truth, these are just renderings.
"""

from __future__ import annotations

import io

from backend.l1arch import service
from backend.planning.exports import MarkdownBlock, _blocks, _diagram_bytes, _filename, _plain


def executive_markdown(project_id: str, l1_element_id: str) -> tuple[bytes, str]:
    summary = service.executive_summary(project_id, l1_element_id)
    title = summary["markdown"].splitlines()[0].lstrip("# ").strip() or "L1 baseline"
    return summary["markdown"].encode("utf-8"), _filename(title, "md")


def _title(project_id: str, l1_element_id: str) -> tuple[str, str]:
    summary = service.executive_summary(project_id, l1_element_id)
    markdown = summary["markdown"]
    title = markdown.splitlines()[0].lstrip("# ").strip() or "L1 baseline"
    return title, markdown


def executive_docx(project_id: str, l1_element_id: str, diagram_images: list[str] | None = None) -> tuple[bytes, str]:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor

    title, markdown = _title(project_id, l1_element_id)
    rendered = _diagram_bytes(diagram_images or [])
    diagram_index = 0
    document = Document()
    document.styles["Normal"].font.name = "Aptos"
    document.styles["Normal"].font.size = Pt(10.5)
    heading = document.add_heading(title, 0)
    heading.runs[0].font.color.rgb = RGBColor(11, 87, 208)

    for block in _blocks(markdown):
        if block.kind == "heading":
            if block.level == 1:
                continue  # already the doc title
            document.add_heading(_plain(block.text), min(block.level, 4))
        elif block.kind in {"bullet", "number"}:
            document.add_paragraph(_plain(block.text), style="List Bullet" if block.kind == "bullet" else "List Number")
        elif block.kind == "quote":
            document.add_paragraph(_plain(block.text), style="Quote")
        elif block.kind in {"code", "mermaid"}:
            if block.kind == "mermaid" and diagram_index < len(rendered) and rendered[diagram_index]:
                document.add_picture(io.BytesIO(rendered[diagram_index]), width=Inches(6.2))
            if block.kind == "mermaid":
                diagram_index += 1
        elif block.kind == "table":
            document.add_paragraph(_plain(block.text))
        else:
            document.add_paragraph(_plain(block.text))
    output = io.BytesIO()
    document.save(output)
    return output.getvalue(), _filename(title, "docx")


def executive_pptx(project_id: str, l1_element_id: str, diagram_images: list[str] | None = None) -> tuple[bytes, str]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    title, markdown = _title(project_id, l1_element_id)
    rendered = _diagram_bytes(diagram_images or [])
    diagram_index = 0
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "L1 Architecture Baseline · Executive summary"

    current_title = "Overview"
    current: list[MarkdownBlock] = []

    def flush() -> None:
        nonlocal current, current_title
        if not current:
            return
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = _plain(current_title)
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, block in enumerate(current):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = _plain(block.text)
            paragraph.level = 1 if block.kind in {"bullet", "number"} else 0
            paragraph.font.size = Pt(16)
        current = []

    for block in _blocks(markdown):
        if block.kind == "heading" and block.level <= 2:
            flush()
            current_title = _plain(block.text)
        elif block.kind == "mermaid":
            flush()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = current_title
            if diagram_index < len(rendered) and rendered[diagram_index]:
                slide.shapes.add_picture(io.BytesIO(rendered[diagram_index]), Inches(1.2), Inches(1.4), width=Inches(10.9))
            diagram_index += 1
        elif block.kind not in {"code"}:
            current.append(block)
    flush()

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue(), _filename(title, "pptx")
