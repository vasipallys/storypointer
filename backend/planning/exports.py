"""Generate portable Word and PowerPoint requirement artifacts."""

from __future__ import annotations

import io
import base64
import binascii
import re
from dataclasses import dataclass

from backend.planning.requirements import get_document
from backend.planning.store import PlanningValidationError


@dataclass
class MarkdownBlock:
    kind: str
    text: str
    level: int = 0


def _blocks(markdown: str) -> list[MarkdownBlock]:
    blocks: list[MarkdownBlock] = []
    paragraph: list[str] = []
    code: list[str] = []
    code_language = ""
    in_code = False

    def flush_paragraph() -> None:
        if paragraph:
            blocks.append(MarkdownBlock("paragraph", " ".join(part.strip() for part in paragraph)))
            paragraph.clear()

    for line in markdown.splitlines():
        fence = re.match(r"^\s*```([\w-]*)\s*$", line)
        if fence:
            if in_code:
                blocks.append(MarkdownBlock("mermaid" if code_language == "mermaid" else "code", "\n".join(code)))
                code.clear()
                code_language = ""
                in_code = False
            else:
                flush_paragraph()
                in_code = True
                code_language = fence.group(1).lower()
            continue
        if in_code:
            code.append(line)
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", line)
        bullet = re.match(r"^\s*[-*+]\s+(.+)$", line)
        numbered = re.match(r"^\s*\d+[.)]\s+(.+)$", line)
        if heading:
            flush_paragraph()
            blocks.append(MarkdownBlock("heading", heading.group(2).strip(), len(heading.group(1))))
        elif bullet:
            flush_paragraph()
            blocks.append(MarkdownBlock("bullet", bullet.group(1).strip()))
        elif numbered:
            flush_paragraph()
            blocks.append(MarkdownBlock("number", numbered.group(1).strip()))
        elif line.startswith(">"):
            flush_paragraph()
            blocks.append(MarkdownBlock("quote", line.lstrip("> ").strip()))
        elif line.strip():
            paragraph.append(line)
        else:
            flush_paragraph()
    flush_paragraph()
    if code:
        blocks.append(MarkdownBlock("mermaid" if code_language == "mermaid" else "code", "\n".join(code)))
    return blocks


def _plain(text: str) -> str:
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"(\*\*|__)(.*?)\1", r"\2", text)
    text = re.sub(r"(\*|_)(.*?)\1", r"\2", text)
    return text.replace("`", "")


def _diagram_bytes(images: list[str]) -> list[bytes | None]:
    from PIL import Image

    decoded: list[bytes | None] = []
    total = 0
    for image in images:
        if not image:
            decoded.append(None)
            continue
        try:
            prefix, encoded = image.split(",", 1)
            if "image/png" not in prefix:
                raise ValueError
            content = base64.b64decode(encoded, validate=True)
            total += len(content)
            if not content.startswith(b"\x89PNG\r\n\x1a\n") or total > 20_000_000:
                raise ValueError
            with Image.open(io.BytesIO(content)) as rendered:
                rendered.verify()
            decoded.append(content)
        except (ValueError, TypeError, binascii.Error, OSError):
            raise PlanningValidationError("Rendered Mermaid images must be valid PNG data URLs") from None
    return decoded


def word_export(
    project_id: str,
    document_id: str,
    diagram_images: list[str] | None = None,
) -> tuple[bytes, str]:
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.shared import Inches, Pt, RGBColor

    requirement = get_document(project_id, document_id)
    rendered_diagrams = _diagram_bytes(diagram_images or [])
    diagram_index = 0
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)
    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5)

    title = document.add_heading(requirement["title"], 0)
    title.runs[0].font.color.rgb = RGBColor(11, 87, 208)
    meta = document.add_paragraph()
    meta.add_run(
        f"L1 requirements  •  Version {requirement['version']}  •  "
        f"{requirement['status'].replace('_', ' ').title()}"
    ).italic = True
    if requirement.get("approved_by"):
        document.add_paragraph(
            f"Approved by {requirement['approved_by']} on {requirement['approved_at']}",
            style="Caption",
        )
    document.add_paragraph("")

    for block in _blocks(requirement["content"]):
        if block.kind == "heading":
            document.add_heading(_plain(block.text), min(block.level, 4))
        elif block.kind in {"bullet", "number"}:
            document.add_paragraph(_plain(block.text), style="List Bullet" if block.kind == "bullet" else "List Number")
        elif block.kind == "quote":
            paragraph = document.add_paragraph(_plain(block.text), style="Quote")
            paragraph.paragraph_format.left_indent = Inches(0.25)
        elif block.kind in {"code", "mermaid"}:
            if block.kind == "mermaid":
                document.add_heading("Mermaid diagram", 3)
                if diagram_index < len(rendered_diagrams) and rendered_diagrams[diagram_index]:
                    document.add_picture(io.BytesIO(rendered_diagrams[diagram_index]), width=Inches(6.4))
                document.add_paragraph("Editable Mermaid definition", style="Caption")
                diagram_index += 1
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.left_indent = Inches(0.2)
            run = paragraph.add_run(block.text)
            run.font.name = "Consolas"
            run.font.size = Pt(8.5)
            run.font.color.rgb = RGBColor(31, 31, 31)
        else:
            document.add_paragraph(_plain(block.text))

    document.add_section(WD_SECTION.NEW_PAGE)
    document.add_heading("Review record", 1)
    document.add_paragraph(
        f"Exported from the audit-preserved requirements workspace at version {requirement['version']}."
    )
    for comment in requirement["comments"]:
        paragraph = document.add_paragraph(style="List Bullet")
        paragraph.add_run(f"{comment['author']} — {comment['status'].title()}: ").bold = True
        paragraph.add_run(comment["body"])
    output = io.BytesIO()
    document.save(output)
    return output.getvalue(), _filename(requirement["title"], "docx")


def powerpoint_export(
    project_id: str,
    document_id: str,
    diagram_images: list[str] | None = None,
) -> tuple[bytes, str]:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    requirement = get_document(project_id, document_id)
    rendered_diagrams = _diagram_bytes(diagram_images or [])
    diagram_index = 0
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = requirement["title"]
    title_slide.placeholders[1].text = (
        f"L1 requirements • Version {requirement['version']} • "
        f"{requirement['status'].replace('_', ' ').title()}"
    )

    current_title = "Requirements overview"
    current: list[MarkdownBlock] = []

    def add_content_slide(title: str, items: list[MarkdownBlock]) -> None:
        if not items:
            return
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = _plain(title)
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, block in enumerate(items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            prefix = "Mermaid diagram definition:\n" if block.kind == "mermaid" else ""
            paragraph.text = prefix + _plain(block.text)
            paragraph.level = 1 if block.kind in {"bullet", "number"} else 0
            paragraph.font.name = "Consolas" if block.kind in {"code", "mermaid"} else "Aptos"
            paragraph.font.size = Pt(12 if block.kind in {"code", "mermaid"} else 18)
            paragraph.space_after = Pt(8)
        footer = slide.shapes.add_textbox(Inches(10.2), Inches(7.05), Inches(2.7), Inches(0.25))
        footer.text_frame.paragraphs[0].text = f"Version {requirement['version']} • {requirement['status']}"
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        footer.text_frame.paragraphs[0].font.size = Pt(8)

    def add_diagram_slide(title: str, block: MarkdownBlock, image: bytes | None) -> None:
        from PIL import Image

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = _plain(title)
        if image:
            with Image.open(io.BytesIO(image)) as rendered:
                image_width, image_height = rendered.size
            scale = min(11.3 / image_width, 5.45 / image_height)
            width = image_width * scale
            height = image_height * scale
            slide.shapes.add_picture(
                io.BytesIO(image),
                Inches((13.333 - width) / 2),
                Inches(1.35 + (5.45 - height) / 2),
                width=Inches(width),
                height=Inches(height),
            )
        else:
            source = slide.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.3), Inches(5.2))
            source.text_frame.text = block.text
            source.text_frame.paragraphs[0].font.name = "Consolas"
            source.text_frame.paragraphs[0].font.size = Pt(16)

    for block in _blocks(requirement["content"]):
        if block.kind == "heading" and block.level <= 2:
            add_content_slide(current_title, current)
            current_title, current = block.text, []
        elif block.kind == "mermaid":
            add_content_slide(current_title, current)
            current = []
            image = rendered_diagrams[diagram_index] if diagram_index < len(rendered_diagrams) else None
            add_diagram_slide(f"{current_title} — diagram", block, image)
            diagram_index += 1
            current_title = f"{current_title} (continued)"
        else:
            current.append(block)
            if len(current) >= 6 or sum(len(item.text) for item in current) > 850:
                add_content_slide(current_title, current)
                current_title = f"{current_title} (continued)"
                current = []
    add_content_slide(current_title, current)

    if requirement["comments"]:
        review_items = [
            MarkdownBlock("bullet", f"{item['author']} ({item['status']}): {item['body']}")
            for item in requirement["comments"][:10]
        ]
        add_content_slide("Review comments", review_items)

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue(), _filename(requirement["title"], "pptx")


def _filename(title: str, extension: str) -> str:
    stem = re.sub(r"[^A-Za-z0-9_-]+", "-", title).strip("-").lower() or "l1-requirements"
    return f"{stem}.{extension}"
